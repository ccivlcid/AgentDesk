"""AgentDesk 프로젝트 소개 — 경영진 보고 PPT 생성"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──────────────────────────────────────────────
BG_DARK = RGBColor(0x0D, 0x11, 0x17)
BG_CARD = RGBColor(0x16, 0x1B, 0x22)
BG_ACCENT = RGBColor(0x1C, 0x22, 0x2B)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
WHITE = RGBColor(0xE6, 0xED, 0xF3)
MUTED = RGBColor(0x8B, 0x94, 0x9E)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
RED = RGBColor(0xEF, 0x44, 0x44)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
CYAN = RGBColor(0x06, 0xB6, 0xD4)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_slide_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=18, color=WHITE,
             bold=False, alignment=PP_ALIGN.LEFT, font_name="Malgun Gothic"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline(slide, left, top, width, height, lines, size=16,
                  font_name="Malgun Gothic"):
    """lines: list of (text, color, bold) tuples"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font_name
    return txBox


def add_accent_bar(slide, left, top, height, color=AMBER):
    return add_shape(slide, left, top, Pt(3), height, fill_color=color)


def add_card(slide, x, y, w, h, title, items, accent=AMBER, title_size=20, item_size=14):
    add_shape(slide, x, y, w, h, fill_color=BG_CARD)
    add_accent_bar(slide, x, y, h, accent)
    add_text(slide, x + Inches(0.3), y + Inches(0.15), w - Inches(0.4), Inches(0.5),
             title, size=title_size, color=accent, bold=True, font_name="Consolas")
    lines = []
    for item in items:
        if isinstance(item, tuple):
            lines.append((f"  {item[0]}  —  {item[1]}", WHITE, False))
        else:
            lines.append((f"  {item}", MUTED, False))
    add_multiline(slide, x + Inches(0.3), y + Inches(0.65), w - Inches(0.4),
                  h - Inches(0.8), lines, size=item_size)


def add_stat_box(slide, x, y, number, label, color=AMBER):
    w, h = Inches(2.4), Inches(1.6)
    add_shape(slide, x, y, w, h, fill_color=BG_CARD)
    add_accent_bar(slide, x, y, h, color)
    add_text(slide, x + Inches(0.3), y + Inches(0.2), w, Inches(0.8),
             number, size=36, color=color, bold=True, font_name="Consolas")
    add_text(slide, x + Inches(0.3), y + Inches(1.0), w, Inches(0.4),
             label, size=14, color=MUTED)


# ════════════════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_shape(slide, Inches(1.5), Inches(2.4), Inches(1.0), Pt(3), fill_color=AMBER)

add_text(slide, Inches(1.5), Inches(2.7), Inches(10), Inches(1),
         "▶ AgentDesk", size=52, color=WHITE, bold=True, font_name="Consolas")
add_text(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.8),
         "Project Operating System", size=30, color=AMBER, font_name="Consolas")
add_text(slide, Inches(1.5), Inches(4.7), Inches(10), Inches(0.8),
         "AI 에이전트로 구동되는 프로젝트 관리 플랫폼", size=22, color=MUTED)
add_text(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
         "어떤 팀이든 맞춤 설계하는 프로젝트 운영체제", size=16, color=MUTED,
         font_name="Malgun Gothic")

add_text(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.5),
         "2026  ·  v2.0  ·  Open Source", size=14, color=RGBColor(0x48, 0x50, 0x5A),
         font_name="Consolas")

# ════════════════════════════════════════════════════════
# SLIDE 2: 해결하는 문제
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "AI 시대, 프로젝트 관리의 새로운 과제", size=30, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.0), Inches(11.5), Pt(1), fill_color=BG_ACCENT)

pains = [
    ("기존 방식", "AI 에이전트 시대", MUTED, AMBER),
    ("사람이 직접 코딩, 디자인, 분석 수행", "AI 에이전트가 CLI로 자율 실행", WHITE, GREEN),
    ("PM이 진행 상황을 수동 확인", "실시간 모니터링 + 자동 리포트", WHITE, GREEN),
    ("엑셀/노션에 흩어진 리스크 관리", "프로젝트별 통합 리스크 추적", WHITE, GREEN),
    ("배포 전 품질 게이트 없음", "자동화된 Gate 심사 프로세스", WHITE, GREEN),
    ("프로젝트 유형마다 다른 도구 사용", "카테고리 템플릿으로 통일", WHITE, GREEN),
]

for i, (left_text, right_text, left_color, right_color) in enumerate(pains):
    y = Inches(1.4 + i * 0.85)
    is_header = i == 0
    bg = BG_ACCENT if is_header else (BG_CARD if i % 2 == 0 else None)
    if bg:
        add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(0.7), fill_color=bg)

    add_text(slide, Inches(1.0), y + Inches(0.12), Inches(5.2), Inches(0.5),
             left_text, size=16 if is_header else 15, color=left_color,
             bold=is_header)

    if not is_header:
        add_text(slide, Inches(6.0), y + Inches(0.12), Inches(0.5), Inches(0.5),
                 "→", size=18, color=AMBER, font_name="Consolas")

    add_text(slide, Inches(6.6), y + Inches(0.12), Inches(5.2), Inches(0.5),
             right_text, size=16 if is_header else 15, color=right_color,
             bold=is_header)

add_text(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5),
         "AgentDesk는 AI 에이전트를 '팀원'으로, 프로젝트를 '운영체제'로 다루는 새로운 패러다임을 제시합니다.",
         size=15, color=AMBER, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE 3: 제품 개요 — What is AgentDesk?
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "AgentDesk란?", size=30, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.0), Inches(11.5), Pt(1), fill_color=BG_ACCENT)

# Left side — definition
add_shape(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(5.3), fill_color=BG_CARD)
add_accent_bar(slide, Inches(0.8), Inches(1.4), Inches(5.3))

add_multiline(slide, Inches(1.2), Inches(1.6), Inches(4.8), Inches(4.8), [
    ("Project Operating System", AMBER, True),
    ("", WHITE, False),
    ("AI 에이전트(Claude, GPT 등)를 CLI 프로세스로", WHITE, False),
    ("실행하며, 프로젝트의 목표·리스크·품질 게이트·", WHITE, False),
    ("결과물을 체계적으로 관리하는 통합 플랫폼", WHITE, False),
    ("", WHITE, False),
    ("▸ 에이전트 ≠ 챗봇", AMBER, True),
    ("  백그라운드 CLI 프로세스로 자율 실행", MUTED, False),
    ("", WHITE, False),
    ("▸ 프로젝트 ≠ 태스크 리스트", AMBER, True),
    ("  목표-리스크-게이트-결과물 4축 관리", MUTED, False),
    ("", WHITE, False),
    ("▸ 카테고리 ≠ 폴더", AMBER, True),
    ("  프로젝트 유형별 KPI/리스크/게이트 스키마", MUTED, False),
], size=15)

# Right side — key numbers
add_stat_box(slide, Inches(6.8), Inches(1.4), "6+", "지원 AI 프로바이더", AMBER)
add_stat_box(slide, Inches(9.5), Inches(1.4), "4축", "프로젝트 거버넌스", GREEN)
add_stat_box(slide, Inches(6.8), Inches(3.3), "6종", "프로젝트 카테고리 템플릿", BLUE)
add_stat_box(slide, Inches(9.5), Inches(3.3), "실시간", "에이전트 모니터링", CYAN)

add_text(slide, Inches(6.8), Inches(5.3), Inches(5.2), Inches(1.2),
         "IT 개발 · 투자 운영 · 리서치 · 마케팅\n어떤 프로젝트 유형이든 템플릿으로 즉시 시작",
         size=14, color=MUTED, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE 4: 핵심 기능 — 4-Quadrant Dashboard
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "핵심 기능 ① — 프로젝트 대시보드 (4-Quadrant)", size=28, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.0), Inches(11.5), Pt(1), fill_color=BG_ACCENT)

quads = [
    ("GOALS", "프로젝트 목표", [
        "핵심 목표 설정 및 진행률 추적",
        "목표별 상태 관리 (active / completed)",
        "프로젝트 성공 기준 정의",
    ], GREEN),
    ("RISKS", "리스크 관리", [
        "위험 요소 식별 및 심각도 분류",
        "완화 계획 수립 + 담당자 지정",
        "자동 리스크 홀드 (심각도 기반)",
    ], RED),
    ("GATES", "품질 게이트", [
        "단계별 검토 체크포인트 설정",
        "PENDING → PASSED / FAILED 심사",
        "게이트 미통과 시 다음 단계 차단",
    ], BLUE),
    ("OUTPUTS", "결과물 추적", [
        "산출물 정의 및 버전 관리",
        "완료 상태 추적 (TODO → DONE)",
        "재사용 가능한 결과물 제안",
    ], AMBER),
]

for i, (code, title, items, color) in enumerate(quads):
    x = Inches(0.8 + i * 3.1)
    h = Inches(4.8)
    add_shape(slide, x, Inches(1.4), Inches(2.8), h, fill_color=BG_CARD)
    add_accent_bar(slide, x, Inches(1.4), h, color)

    add_text(slide, x + Inches(0.3), Inches(1.6), Inches(2.4), Inches(0.4),
             code, size=22, color=color, bold=True, font_name="Consolas")
    add_text(slide, x + Inches(0.3), Inches(2.1), Inches(2.4), Inches(0.4),
             title, size=16, color=WHITE, bold=True)

    lines = [(f"  {item}", MUTED, False) for item in items]
    add_multiline(slide, x + Inches(0.3), Inches(2.7), Inches(2.4), Inches(3.2),
                  lines, size=13)

add_text(slide, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.4),
         "모든 프로젝트는 이 4축으로 건강 상태를 한눈에 파악할 수 있습니다",
         size=14, color=AMBER, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE 5: 핵심 기능 — AI 에이전트 오케스트레이션
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         "핵심 기능 ② — AI 에이전트 오케스트레이션", size=28, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.0), Inches(11.5), Pt(1), fill_color=BG_ACCENT)

# Flow diagram — horizontal steps
steps = [
    ("01", "에이전트 생성", "이름, 역할, 사고방식\n(페르소나) 설정", AMBER),
    ("02", "CLI 프로바이더", "Claude / GPT /\nGemini / Ollama 선택", BLUE),
    ("03", "프로젝트 배정", "팀 구성 + 부서 배치\n권한 설정", GREEN),
    ("04", "태스크 할당", "칸반 보드에서\n드래그 앤 드롭", CYAN),
    ("05", "자율 실행", "CLI 프로세스로\n백그라운드 실행", GREEN),
    ("06", "모니터링", "실시간 터미널 로그\n+ Heartbeat 상태", AMBER),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(0.5 + i * 2.1)
    y = Inches(1.5)
    w = Inches(1.9)
    h = Inches(2.5)

    add_shape(slide, x, y, w, h, fill_color=BG_CARD)
    add_accent_bar(slide, x, y, h, color)

    add_text(slide, x + Inches(0.2), y + Inches(0.15), w, Inches(0.4),
             num, size=24, color=color, bold=True, font_name="Consolas")
    add_text(slide, x + Inches(0.2), y + Inches(0.6), w - Inches(0.3), Inches(0.4),
             title, size=14, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(1.1), w - Inches(0.3), Inches(1.2),
             desc, size=12, color=MUTED)

# Execution states
add_text(slide, Inches(0.8), Inches(4.3), Inches(11.5), Inches(0.4),
         "실행 상태 흐름", size=18, color=WHITE, bold=True)

state_flow = "queued  →  claiming  →  workspace_preparing  →  ready  →  running  →  awaiting_review  →  succeeded / failed"
add_shape(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.6), fill_color=BG_CARD)
add_text(slide, Inches(1.0), Inches(4.85), Inches(11.0), Inches(0.5),
         state_flow, size=13, color=GREEN, font_name="Consolas", alignment=PP_ALIGN.CENTER)

# Bottom — supporting systems
supports = [
    ("SKILLS", "에이전트 역량 라이브러리"),
    ("RULES", "행동 제약 규칙"),
    ("MEMORY", "장기 컨텍스트 저장"),
    ("HOOKS", "이벤트 기반 자동화"),
]
for i, (label, desc) in enumerate(supports):
    x = Inches(0.8 + i * 3.1)
    add_shape(slide, x, Inches(5.7), Inches(2.8), Inches(1.0), fill_color=BG_CARD)
    add_text(slide, x + Inches(0.2), Inches(5.8), Inches(2.4), Inches(0.4),
             label, size=14, color=AMBER, bold=True, font_name="Consolas")
    add_text(slide, x + Inches(0.2), Inches(6.2), Inches(2.4), Inches(0.4),
             desc, size=12, color=MUTED)

# ════════════════════════════════════════════════════════
# SLIDE 6: 주요 화면 구성
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "주요 화면 구성", size=30, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.0), Inches(11.5), Pt(1), fill_color=BG_ACCENT)

views = [
    ("DASHBOARD", "프로젝트 건강 상태 한눈에", "4-Quadrant 패널\n목표 진행률 · 리스크 현황\n게이트 심사 · 결과물 추적"),
    ("WORKMAP", "에이전트 실시간 활동 모니터", "부서별 에이전트 상태\n실행 중 태스크 표시\n프로세스 모니터링"),
    ("TASK BOARD", "칸반 기반 태스크 관리", "상태별 컬럼 (6단계)\n필터 (부서/유형/에이전트)\nGantt · DAG 뷰"),
    ("HEARTBEAT", "에이전트 프로세스 모니터", "Uptime · CPU 사용량\n실행 로그 뷰어\n이상 감지 알림"),
    ("LIBRARY", "지식 자산 관리", "Skills · Rules\nMemory · Hooks\n에이전트별/전역 관리"),
    ("SETTINGS", "시스템 설정", "OAuth 연동\nAPI 키 관리\n카테고리 에디터"),
]

for i, (title, subtitle, desc) in enumerate(views):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.0)
    y = Inches(1.4 + row * 2.9)
    w, h = Inches(3.6), Inches(2.6)

    add_shape(slide, x, y, w, h, fill_color=BG_CARD)
    add_accent_bar(slide, x, y, h)

    add_text(slide, x + Inches(0.3), y + Inches(0.15), w, Inches(0.4),
             title, size=18, color=AMBER, bold=True, font_name="Consolas")
    add_text(slide, x + Inches(0.3), y + Inches(0.55), w - Inches(0.4), Inches(0.4),
             subtitle, size=14, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(1.0), w - Inches(0.4), Inches(1.5),
             desc, size=13, color=MUTED)

# ════════════════════════════════════════════════════════
# SLIDE 7: 기술 스택
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "기술 스택", size=30, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.0), Inches(11.5), Pt(1), fill_color=BG_ACCENT)

# Left — Frontend
add_card(slide, Inches(0.8), Inches(1.4), Inches(3.6), Inches(5.0),
         "FRONTEND", [
             ("React", "19.2 + TypeScript 5.9"),
             ("Vite", "7.2 — 빌드 도구"),
             ("Tailwind CSS", "4.1 — 스타일링"),
             ("Framer Motion", "12 — 애니메이션"),
             ("Lucide", "아이콘 시스템"),
             ("WebSocket", "실시간 통신"),
         ], BLUE)

# Center — Backend
add_card(slide, Inches(4.8), Inches(1.4), Inches(3.6), Inches(5.0),
         "BACKEND", [
             ("Node.js", "≥ 22 런타임"),
             ("Express", "5.2 — 웹 프레임워크"),
             ("SQLite", "임베디드 DB"),
             ("WebSocket", "ws — 실시간 업데이트"),
             ("Sharp", "이미지 처리"),
             ("20+ 테이블", "스키마 관리"),
         ], GREEN)

# Right — Infrastructure
add_card(slide, Inches(8.8), Inches(1.4), Inches(3.6), Inches(5.0),
         "INFRA & AI", [
             ("Electron", "데스크톱 앱"),
             ("Anthropic", "Claude Code CLI"),
             ("OpenAI", "GPT CLI"),
             ("Google", "Gemini"),
             ("Ollama", "로컬 LLM"),
             ("Cursor", "IDE 연동"),
         ], AMBER)

add_text(slide, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.4),
         "218+ 컴포넌트  ·  20+ DB 테이블  ·  6개 CLI 프로바이더  ·  Open Source",
         size=15, color=AMBER, alignment=PP_ALIGN.CENTER, font_name="Consolas")

# ════════════════════════════════════════════════════════
# SLIDE 8: 카테고리 시스템
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         "핵심 기능 ③ — 카테고리 템플릿 시스템", size=28, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.0), Inches(11.5), Pt(1), fill_color=BG_ACCENT)

add_text(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.5),
         "프로젝트 유형별로 KPI · 리스크 · 게이트 · 결과물 스키마를 사전 정의하고 재사용",
         size=16, color=MUTED, alignment=PP_ALIGN.CENTER)

categories = [
    ("IT DELIVERY", "소프트웨어 개발 프로젝트\n코드 품질, 테스트 커버리지\n코드 리뷰 게이트", BLUE),
    ("INVESTMENT", "투자 운영 프로젝트\n수익률, 리스크 지표\n실사 검토 게이트", GREEN),
    ("RESEARCH", "리서치/전략 프로젝트\n논문 수, 인사이트 KPI\n피어 리뷰 게이트", CYAN),
    ("MARKETING", "마케팅/그로스 프로젝트\nCAC, LTV, 전환율\n캠페인 승인 게이트", RED),
    ("CUSTOM", "사용자 정의 템플릿\n자유 스키마 구성\n팀 맞춤 워크플로우", AMBER),
]

for i, (title, desc, color) in enumerate(categories):
    x = Inches(0.5 + i * 2.5)
    y = Inches(2.1)
    w, h = Inches(2.3), Inches(2.8)

    add_shape(slide, x, y, w, h, fill_color=BG_CARD)
    add_accent_bar(slide, x, y, h, color)

    add_text(slide, x + Inches(0.25), y + Inches(0.15), w - Inches(0.3), Inches(0.5),
             title, size=14, color=color, bold=True, font_name="Consolas")
    add_text(slide, x + Inches(0.25), y + Inches(0.7), w - Inches(0.3), Inches(2.0),
             desc, size=12, color=MUTED)

# Version locking explanation
add_shape(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(1.5), fill_color=BG_CARD)
add_accent_bar(slide, Inches(0.8), Inches(5.3), Inches(1.5))

add_multiline(slide, Inches(1.2), Inches(5.4), Inches(10.5), Inches(1.3), [
    ("버전 잠금(Version Locking) 아키텍처", AMBER, True),
    ("프로젝트 생성 시 카테고리 버전이 고정됩니다. 이후 카테고리가 업데이트되더라도 기존 프로젝트에는 영향 없음.", MUTED, False),
    ("→ 과거 프로젝트의 데이터 무결성 보장  ·  신규 프로젝트만 최신 템플릿 적용  ·  필요 시 opt-in 업그레이드", WHITE, False),
], size=14)

# ════════════════════════════════════════════════════════
# SLIDE 9: 디자인 아이덴티티
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "디자인 아이덴티티", size=30, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.0), Inches(11.5), Pt(1), fill_color=BG_ACCENT)

add_text(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.5),
         "Professional Tool  ×  Terminal Precision  —  90% 프로페셔널, 10% 터미널",
         size=16, color=AMBER, alignment=PP_ALIGN.CENTER, font_name="Consolas")

design_items = [
    ("THEME", "다크 모드 기본 (#0D1117)", "전문 도구의 집중력 + 눈의 피로 최소화"),
    ("ACCENT", "Amber (#F59E0B)", "브랜드 컬러, CTA 버튼, 실행 상태 표시"),
    ("TYPOGRAPHY", "Sans-serif (UI) + Mono (데이터)", "전문성과 정밀함의 시각적 분리"),
    ("SHAPE", "border-radius: 0 (직각)", "터미널 네이티브 미학, 군더더기 없는 인터페이스"),
    ("ANIMATION", "부트 시퀀스 + 상태 펄스", "터미널 기동 연출, 실행 중 에이전트 시각 피드백"),
    ("LAYOUT", "리스트/테이블 우선", "높은 정보 밀도, Paperclip 스타일 효율"),
]

for i, (label, value, desc) in enumerate(design_items):
    y = Inches(1.9 + i * 0.85)
    add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(0.72), fill_color=BG_CARD if i % 2 == 0 else None)
    add_accent_bar(slide, Inches(0.8), y, Inches(0.72))

    add_text(slide, Inches(1.2), y + Inches(0.08), Inches(1.8), Inches(0.35),
             label, size=14, color=AMBER, bold=True, font_name="Consolas")
    add_text(slide, Inches(3.2), y + Inches(0.08), Inches(4), Inches(0.35),
             value, size=14, color=WHITE, bold=True)
    add_text(slide, Inches(7.5), y + Inches(0.08), Inches(4.5), Inches(0.35),
             desc, size=13, color=MUTED)

# ════════════════════════════════════════════════════════
# SLIDE 10: 로드맵 & 비전
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "개발 현황 & 비전", size=30, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.0), Inches(11.5), Pt(1), fill_color=BG_ACCENT)

# Completed phases
add_text(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.4),
         "완료된 개발 단계", size=18, color=GREEN, bold=True)

completed = [
    ("Phase 1-2", "DB 스키마 · 카테고리 시스템 · 대시보드 2.0"),
    ("Phase 3", "카테고리 에디터 · 자동화 (KPI, 리스크, 결과물)"),
    ("Phase 4", "TaskBoard · AgentManager 컨텍스트 통합"),
    ("Phase 5", "오피스팩 제거 · 프로젝트 생성 통일화"),
    ("Phase 6", "전체 화면 프로젝트 컨텍스트 연결"),
]

for i, (phase, desc) in enumerate(completed):
    y = Inches(1.8 + i * 0.55)
    add_text(slide, Inches(1.0), y, Inches(1.5), Inches(0.4),
             f"✓ {phase}", size=13, color=GREEN, font_name="Consolas")
    add_text(slide, Inches(2.8), y, Inches(4), Inches(0.4),
             desc, size=13, color=MUTED)

# Vision / Next
add_text(slide, Inches(6.8), Inches(1.3), Inches(5), Inches(0.4),
         "향후 로드맵", size=18, color=AMBER, bold=True)

vision = [
    ("UI/UX 리뉴얼", "CLI Concept 전면 적용\n디자인 토큰 통일, 피드백 시스템 교체"),
    ("키보드 퍼스트", "Cmd+K 팔레트, 단축키 체계\n파워유저 생산성 극대화"),
    ("URL 라우팅", "딥링크, 북마크, 뒤로가기\n웹 네이티브 네비게이션"),
    ("모바일 반응형", "태블릿/모바일 대응\n현장 모니터링 지원"),
]

for i, (title, desc) in enumerate(vision):
    y = Inches(1.8 + i * 1.2)
    add_shape(slide, Inches(6.8), y, Inches(5.2), Inches(1.0), fill_color=BG_CARD)
    add_accent_bar(slide, Inches(6.8), y, Inches(1.0))
    add_text(slide, Inches(7.2), y + Inches(0.08), Inches(4.5), Inches(0.35),
             title, size=14, color=AMBER, bold=True)
    add_text(slide, Inches(7.2), y + Inches(0.4), Inches(4.5), Inches(0.55),
             desc, size=12, color=MUTED)

# ════════════════════════════════════════════════════════
# SLIDE 11: 핵심 가치 & 마무리
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_shape(slide, Inches(1.5), Inches(1.3), Inches(1.0), Pt(3), fill_color=AMBER)

add_text(slide, Inches(1.5), Inches(1.6), Inches(10), Inches(0.8),
         "▶ AgentDesk", size=40, color=WHITE, bold=True, font_name="Consolas")
add_text(slide, Inches(1.5), Inches(2.4), Inches(10), Inches(0.5),
         "AI 에이전트를 팀원으로, 프로젝트를 운영체제로", size=20, color=AMBER)

values = [
    ("프로젝트 중심", "모든 것이 프로젝트 컨텍스트 안에서 동작\n목표-리스크-게이트-결과물 4축 거버넌스"),
    ("AI 네이티브", "에이전트가 CLI 프로세스로 자율 실행\n6개 AI 프로바이더 지원, 실시간 모니터링"),
    ("템플릿 기반", "카테고리로 프로젝트 유형 표준화\n한번 설계, 반복 사용, 버전 관리"),
    ("오픈소스", "투명한 코드, 커뮤니티 확장 가능\nReact + TypeScript + SQLite"),
]

for i, (title, desc) in enumerate(values):
    x = Inches(1.5 + i * 2.7)
    y = Inches(3.3)
    w, h = Inches(2.4), Inches(2.8)

    add_shape(slide, x, y, w, h, fill_color=BG_CARD)
    add_accent_bar(slide, x, y, h)

    add_text(slide, x + Inches(0.25), y + Inches(0.15), w - Inches(0.3), Inches(0.5),
             title, size=16, color=AMBER, bold=True)
    add_text(slide, x + Inches(0.25), y + Inches(0.65), w - Inches(0.3), Inches(2.0),
             desc, size=13, color=MUTED)

add_text(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.5),
         "github.com/agentdesk  ·  Open Source  ·  2026",
         size=14, color=RGBColor(0x48, 0x50, 0x5A), font_name="Consolas",
         alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════
output_path = r"C:\PythonProjects\AgentDesk\docs\design\AgentDesk_프로젝트소개_경영진보고.pptx"
prs.save(output_path)
print(f"PPT saved: {output_path}")
