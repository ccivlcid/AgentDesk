"""AgentDesk 2.0 화면 리디자인 — 경영진 보고 PPT 생성"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──────────────────────────────────────────────
BG_DARK = RGBColor(0x0D, 0x11, 0x17)
BG_CARD = RGBColor(0x16, 0x1B, 0x22)
BG_ACCENT = RGBColor(0x1C, 0x22, 0x2B)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
WHITE = RGBColor(0xE6, 0xED, 0xF3)
MUTED = RGBColor(0x8B, 0x94, 0x9E)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
RED = RGBColor(0xEF, 0x44, 0x44)
BLUE = RGBColor(0x3B, 0x82, 0xF6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=18, color=WHITE,
             bold=False, alignment=PP_ALIGN.LEFT, font_name="Malgun Gothic"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, size=16,
                    color=WHITE, font_name="Malgun Gothic"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)

        # Support (label, value) tuples or plain strings
        if isinstance(item, tuple):
            run1 = p.add_run()
            run1.text = item[0]
            run1.font.size = Pt(size)
            run1.font.color.rgb = AMBER
            run1.font.bold = True
            run1.font.name = font_name

            run2 = p.add_run()
            run2.text = f"  {item[1]}"
            run2.font.size = Pt(size)
            run2.font.color.rgb = color
            run2.font.name = font_name
        else:
            run = p.add_run()
            run.text = f"  {item}"
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.name = font_name
    return txBox


def add_accent_bar(slide, left, top, width=Pt(3), height=Inches(0.5)):
    return add_shape(slide, left, top, width, height, fill_color=AMBER)


def add_tag(slide, left, top, text, color=AMBER):
    w, h = Inches(1.8), Inches(0.35)
    shape = add_shape(slide, left, top, w, h, line_color=color)
    shape.text_frame.paragraphs[0].text = text
    shape.text_frame.paragraphs[0].font.size = Pt(11)
    shape.text_frame.paragraphs[0].font.color.rgb = color
    shape.text_frame.paragraphs[0].font.name = "Consolas"
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape


# ════════════════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

# Amber accent line
add_shape(slide, Inches(1.5), Inches(2.8), Inches(0.8), Pt(3), fill_color=AMBER)

add_text(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(1),
         "AgentDesk 2.0", size=44, color=WHITE, bold=True, font_name="Consolas")
add_text(slide, Inches(1.5), Inches(3.9), Inches(10), Inches(0.8),
         "전체 화면 리디자인 계획", size=28, color=MUTED)
add_text(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(0.5),
         "CLI Concept  ·  AI 에이전트 프로세스 관리 전문 터미널", size=16, color=AMBER,
         font_name="Consolas")
add_text(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.5),
         "2026-03-11  ·  v1.0", size=14, color=MUTED, font_name="Consolas")

# ════════════════════════════════════════════════════════
# SLIDE 2: 현재 문제점 요약
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "현재 문제점 요약", size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.0), Inches(11.5), Pt(1), fill_color=BG_ACCENT)

# Problem cards — 3 columns
problems = [
    ("UX 일관성 부재", [
        "border-radius 2px~8px 혼용",
        "폰트 mono/sans 혼용",
        "한/영 레이블 혼재",
        "카드/리스트/테이블 혼용",
    ]),
    ("사용자 피드백 부재", [
        "window.alert() — 21곳",
        "window.confirm() — 9곳",
        ".catch(() => {}) — 12곳",
        "에러 발생 시 무반응",
    ]),
    ("정보 밀도 낮음", [
        "카드 그리드 → 공간 낭비",
        "요약 정보 별도 카드 3개",
        "core_goal 미노출",
        "설정 8탭 평면 나열",
    ]),
]

for i, (title, items) in enumerate(problems):
    x = Inches(0.8 + i * 4.0)
    add_shape(slide, x, Inches(1.4), Inches(3.6), Inches(4.5), fill_color=BG_CARD)
    add_accent_bar(slide, x, Inches(1.4), Pt(3), Inches(4.5))
    add_text(slide, x + Inches(0.3), Inches(1.6), Inches(3.2), Inches(0.5),
             title, size=20, color=AMBER, bold=True)
    add_bullet_list(slide, x + Inches(0.3), Inches(2.3), Inches(3.2), Inches(3.5),
                    items, size=15, color=MUTED)

# Bottom stats
add_text(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5),
         "총 42곳의 네이티브 다이얼로그  ·  CSS !important 380곳  ·  rgba() 하드코딩 313곳",
         size=14, color=RED, font_name="Consolas", alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE 3: 디자인 방향
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "디자인 방향: CLI Concept", size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.0), Inches(11.5), Pt(1), fill_color=BG_ACCENT)

principles = [
    ("SHARP", "border-radius: 0 — 모든 UI 요소에 직각 적용 (아바타 제외)"),
    ("MONO", "데이터/상태/ID에 JetBrains Mono, UI 레이블에 system-ui sans-serif"),
    ("STATUS", "● RUNNING  ○ IDLE  ✕ FAILED  ✓ DONE  ⋯ PENDING — 통일 배지 시스템"),
    ("DENSITY", "카드 그리드 → 리스트/테이블 행으로 전환, 정보 밀도 향상"),
    ("FEEDBACK", "Toast + ConfirmDialog로 네이티브 다이얼로그 완전 교체"),
    ("KEYBOARD", "Cmd+K 팔레트, N(새 태스크), /(검색), J/K(이동) — 파워유저 지원"),
]

for i, (label, desc) in enumerate(principles):
    y = Inches(1.4 + i * 0.9)
    add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(0.75), fill_color=BG_CARD)
    add_accent_bar(slide, Inches(0.8), y, Pt(3), Inches(0.75))

    add_text(slide, Inches(1.2), y + Inches(0.12), Inches(1.5), Inches(0.5),
             label, size=18, color=AMBER, bold=True, font_name="Consolas")
    add_text(slide, Inches(2.8), y + Inches(0.12), Inches(9.2), Inches(0.5),
             desc, size=15, color=WHITE)

# ════════════════════════════════════════════════════════
# SLIDE 4: 공통 변경 — 디자인 토큰
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "공통 변경: 디자인 토큰 & 피드백", size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.0), Inches(11.5), Pt(1), fill_color=BG_ACCENT)

# Design tokens table
tokens = [
    ("항목", "현재", "변경 후"),
    ("border-radius", "2px ~ 8px 혼용", "0 (아바타 제외)"),
    ("폰트 (레이블)", "mono/sans 혼용", "system-ui sans-serif"),
    ("폰트 (데이터)", "혼용", "JetBrains Mono"),
    ("hover 배경", "다양", "rgba(255,255,255,0.04)"),
    ("선택 배경", "다양", "rgba(245,158,11,0.08)"),
    ("window.alert()", "21곳", "useToast() → 스택 토스트"),
    ("window.confirm()", "9곳", "useConfirm() → ConfirmDialog"),
    (".catch(() => {})", "12곳", "error toast 표시"),
]

for row_i, (col1, col2, col3) in enumerate(tokens):
    y = Inches(1.3 + row_i * 0.55)
    is_header = row_i == 0
    bg = BG_ACCENT if is_header else (BG_CARD if row_i % 2 == 1 else None)

    if bg:
        add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(0.5), fill_color=bg)

    c1_color = AMBER if is_header else MUTED
    c2_color = AMBER if is_header else RED
    c3_color = AMBER if is_header else GREEN

    add_text(slide, Inches(1.0), y + Inches(0.05), Inches(3), Inches(0.4),
             col1, size=14, color=c1_color, bold=is_header, font_name="Consolas")
    add_text(slide, Inches(4.2), y + Inches(0.05), Inches(3.5), Inches(0.4),
             col2, size=14, color=c2_color, font_name="Consolas")
    add_text(slide, Inches(8.0), y + Inches(0.05), Inches(4), Inches(0.4),
             col3, size=14, color=c3_color, bold=is_header, font_name="Consolas")

# ════════════════════════════════════════════════════════
# SLIDE 5: 주요 화면 변경 (1) — Dashboard, WorkMap, TaskBoard
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "주요 화면 변경 (1/2)", size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.0), Inches(11.5), Pt(1), fill_color=BG_ACCENT)

screens1 = [
    ("DASHBOARD", [
        ("core_goal", "헤더 아래 부제목으로 표시"),
        ("상태 요약", "배지 3개 → 인라인 요약 바"),
        ("탭/패널 레이블", "한/영 혼용 → 영문 대문자 MONO"),
        ("rounded 제거", "모든 요소 border-radius: 0"),
    ]),
    ("WORKMAP", [
        ("뷰 헤더", "프로젝트명 + online/running 카운트"),
        ("필터 탭 추가", "[ALL] [RUNNING] [IDLE] [OFFLINE]"),
        ("부서 그룹화", "부서별 접이식 섹션"),
        ("최소 변경", "이미 가장 CLI스러운 화면"),
    ]),
    ("TASK BOARD", [
        ("컬럼 헤더", "TODO [4] — 태스크 수 배지 추가"),
        ("태스크 ID", "#1234 좌상단 표시 (mono)"),
        ("뷰 전환", "[BOARD] [GANTT] [DAG] bracket 스타일"),
        ("단축키", "N → 새 태스크 생성"),
    ]),
]

for i, (title, items) in enumerate(screens1):
    x = Inches(0.8 + i * 4.0)
    add_shape(slide, x, Inches(1.4), Inches(3.6), Inches(5.0), fill_color=BG_CARD)
    add_accent_bar(slide, x, Inches(1.4), Pt(3), Inches(5.0))

    add_text(slide, x + Inches(0.3), Inches(1.6), Inches(3.2), Inches(0.5),
             title, size=20, color=AMBER, bold=True, font_name="Consolas")

    add_bullet_list(slide, x + Inches(0.3), Inches(2.3), Inches(3.2), Inches(4.0),
                    items, size=14)

# ════════════════════════════════════════════════════════
# SLIDE 6: 주요 화면 변경 (2) — Agent, Heartbeat, Settings
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "주요 화면 변경 (2/2)", size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.0), Inches(11.5), Pt(1), fill_color=BG_ACCENT)

screens2 = [
    ("AGENTS", [
        ("레이아웃", "카드 그리드 → 리스트 행"),
        ("요약 카드 3개", "헤더 인라인: 5 total · 1 running"),
        ("생성 폼", "13필드 한 화면 → 2단계 분리"),
        ("alert() 5곳", "Toast + inline 에러로 교체"),
    ]),
    ("HEARTBEAT", [
        ("Watchlist", "프로세스 테이블 (mono columns)"),
        ("컬럼", "AGENT / STATUS / UPTIME / CPU"),
        ("로그 검색", "/ 키로 vim-style 활성화"),
        ("alert/confirm", "6곳 → Toast + ConfirmDialog"),
    ]),
    ("SETTINGS", [
        ("8 flat tabs", "좌측 사이드바 3그룹으로 재편"),
        ("그룹", "BASIC / INTEGRATIONS / ADVANCED"),
        ("CLI Provider", "ⓘ 툴팁 설명 추가"),
        ("네비게이션", "마우스 → 화살표 키 지원"),
    ]),
]

for i, (title, items) in enumerate(screens2):
    x = Inches(0.8 + i * 4.0)
    add_shape(slide, x, Inches(1.4), Inches(3.6), Inches(5.0), fill_color=BG_CARD)
    add_accent_bar(slide, x, Inches(1.4), Pt(3), Inches(5.0))

    add_text(slide, x + Inches(0.3), Inches(1.6), Inches(3.2), Inches(0.5),
             title, size=20, color=AMBER, bold=True, font_name="Consolas")

    add_bullet_list(slide, x + Inches(0.3), Inches(2.3), Inches(3.2), Inches(4.0),
                    items, size=14)

# ════════════════════════════════════════════════════════
# SLIDE 7: 구현 로드맵
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "구현 로드맵", size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.0), Inches(11.5), Pt(1), fill_color=BG_ACCENT)

phases = [
    ("P0", "즉시", "피드백 시스템 완성", [
        "alert 21곳 → useToast()",
        "confirm 9곳 → useConfirm()",
        ".catch 12곳 → error toast",
    ], GREEN),
    ("P1", "다음 스프린트", "레이아웃 + CLI 스타일", [
        "Dashboard 상태바 + core_goal",
        "카드 → 리스트 전환 (5개 뷰)",
        "Settings 탭 → 사이드바",
        "Agent 생성 폼 2단계 분리",
    ], AMBER),
    ("P2", "백로그", "키보드 + 커맨드", [
        "Cmd+K 커맨드 팔레트",
        "단축키 (N / J / K / /)",
        "WorkMap 필터 + 부서 그룹",
    ], BLUE),
    ("P3", "장기", "구조 개선", [
        "URL 라우팅 (딥링크)",
        "App.tsx 상태 리팩터링",
        "CSS 정리 (380 !important)",
    ], MUTED),
]

for i, (priority, timeline, title, items, color) in enumerate(phases):
    x = Inches(0.8 + i * 3.1)
    card_h = Inches(5.0)
    add_shape(slide, x, Inches(1.4), Inches(2.8), card_h, fill_color=BG_CARD)
    add_accent_bar(slide, x, Inches(1.4), Pt(3), card_h)

    # Priority tag
    add_tag(slide, x + Inches(0.3), Inches(1.6), f"{priority} · {timeline}", color)

    add_text(slide, x + Inches(0.3), Inches(2.15), Inches(2.4), Inches(0.5),
             title, size=16, color=WHITE, bold=True)

    add_bullet_list(slide, x + Inches(0.3), Inches(2.7), Inches(2.4), Inches(3.5),
                    items, size=13, color=MUTED)

# ════════════════════════════════════════════════════════
# SLIDE 8: 기대 효과
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "기대 효과", size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.0), Inches(11.5), Pt(1), fill_color=BG_ACCENT)

effects = [
    ("01", "UX 일관성 확보", "디자인 토큰 통일로 전체 UI의 시각적 통일감 달성\n"
     "CLI Concept 아이덴티티를 모든 화면에 일관 적용"),
    ("02", "사용자 신뢰도 향상", "네이티브 alert/confirm 42곳 완전 제거\n"
     "에러 무시(.catch) 12곳 해소 → 모든 오류에 피드백 제공"),
    ("03", "정보 밀도 향상", "카드 그리드 → 리스트/테이블로 전환 (5개 뷰)\n"
     "동일 화면에서 2~3배 더 많은 정보 노출"),
    ("04", "파워유저 생산성", "Cmd+K 팔레트, 단축키 체계 도입\n"
     "키보드만으로 핵심 워크플로우 수행 가능"),
]

for i, (num, title, desc) in enumerate(effects):
    y = Inches(1.4 + i * 1.4)
    add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(1.2), fill_color=BG_CARD)
    add_accent_bar(slide, Inches(0.8), y, Pt(3), Inches(1.2))

    add_text(slide, Inches(1.2), y + Inches(0.1), Inches(0.7), Inches(0.5),
             num, size=28, color=AMBER, bold=True, font_name="Consolas")
    add_text(slide, Inches(2.0), y + Inches(0.1), Inches(3), Inches(0.5),
             title, size=20, color=WHITE, bold=True)
    add_text(slide, Inches(2.0), y + Inches(0.55), Inches(9.5), Inches(0.6),
             desc, size=14, color=MUTED)

# ════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════
output_path = r"C:\PythonProjects\AgentDesk\docs\design\AgentDesk_2.0_리디자인_경영진보고.pptx"
prs.save(output_path)
print(f"PPT saved: {output_path}")
